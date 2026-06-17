import glob, pickle
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Inches, Pt
from PIL import Image

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

def px_to_inches(path):
    im = Image.open(path)
    width = im.width / im.info['dpi'][0]
    height = im.height / im.info['dpi'][1]
    return (width, height)

def slides(jpegs, slide_name, jpeg_names=None, legend_jpegs=None):

    prs = Presentation()
    slide_size = (16, 9)
    prs.slide_width, prs.slide_height = Inches(slide_size[0]), Inches(slide_size[1])
    
    for i_index,i_jpeg in enumerate(jpegs):

        #print('slides: i_index=', i_index)

        img = px_to_inches(i_jpeg)

        if not legend_jpegs is None:
            if len(legend_jpegs[i_index])>0:
                img_legend=px_to_inches(legend_jpegs[i_index])

        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        img_h=img[1]; img_w=img[0]

        slide_h=slide_size[1]
        slide_w=slide_h*img[0]/img[1]

        if slide_w>slide_size[0]:
            slide_h=slide_h*slide_size[0]/slide_w
            slide_w=slide_size[0]

        if slide_h>slide_size[1]:
            slide_w=slide_w*slide_size[1]/slide_h
            slide_h=slide_size[1]

        left=slide_size[0]-slide_w
        top=slide_size[1]-slide_h

        left=Inches(left); top=Inches(top)

#        pic = slide.shapes.add_picture(i_jpeg, left, top, width=Inches(img[0]*2), height=Inches(img[1]*2))
        pic = slide.shapes.add_picture(i_jpeg, left, top, width=Inches(slide_w), height=Inches(slide_h))

        if not legend_jpegs is None:
            if len(legend_jpegs[i_index])>0:           
                leg_top=slide_size[1]-slide_h*0.6
                leg_top=Inches(leg_top)

                pic = slide.shapes.add_picture(legend_jpegs[i_index], left*0.2, leg_top, width=Inches(slide_w)*0.2, height=Inches(slide_h)*0.5)
        
        if jpeg_names is not None:
            t_left = Inches(0.01)
            t_top = Inches(.01)
            t_width = Inches(2.0)
            t_height = Inches(0.4)
            txBox = slide.shapes.add_textbox(t_left, t_top, t_width,t_height)

            tf = txBox.text_frame.paragraphs[0]
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.word_wrap = False
            tf.margin_top = 0
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            run = tf.add_run()


            if isinstance(jpeg_names[i_index], list):
                run.text = jpeg_names[i_index][0]
            else:
                run.text = jpeg_names[i_index]

            font = run.font
            font.name = 'Calibri'
            font.size = Pt(28)
            font.bold = True
            font.italic = None  # cause value to be inherited from theme
            font.color.theme_color = MSO_THEME_COLOR.ACCENT_6


    prs.save(slide_name)


def extract_planet_date(filename):
    tmpvar=filename.split('/')[-1]
    iloc=tmpvar.find('_3B')
    return tmpvar[0:iloc]

if __name__=="__main__":
    
    jpegs=['soft_max_colorbar_only.jpeg','nn_colorbar.jpeg']*2
    legend_jpegs=['','','soft_max_colorbar_only.jpeg','nn_colorbar.jpeg']

    jpeg_names=['a','b']*2

    slide_name='test.pptx'
    slides(jpegs, slide_name, jpeg_names=jpeg_names,legend_jpegs=legend_jpegs)
